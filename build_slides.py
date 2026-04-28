"""Build the CS 2640 presentation as PowerPoint."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

OUT = "/n/home03/minli/caching_competition/CS2640_LearnedCaching.pptx"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]

NAVY = RGBColor(0x14, 0x2A, 0x4E)
ACCENT = RGBColor(0xC0, 0x39, 0x2B)
MUTED = RGBColor(0x55, 0x55, 0x55)
LIGHT = RGBColor(0xF4, 0xF4, 0xF4)
GREEN = RGBColor(0x1E, 0x88, 0x4A)


def add_text(slide, left, top, width, height, text, *, size=18, bold=False,
             color=NAVY, align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.05)
    if isinstance(text, str):
        text = [text]
    for i, line in enumerate(text):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return tb


def add_bullets(slide, left, top, width, height, items, *, size=16, color=NAVY):
    """items: list of (text, indent_level) tuples or plain strings."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if isinstance(item, tuple):
            text, lvl = item
        else:
            text, lvl = item, 0
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = lvl
        p.alignment = PP_ALIGN.LEFT
        bullet = "  " * lvl + ("•  " if lvl == 0 else "–  ")
        run = p.add_run()
        run.text = bullet + text
        run.font.size = Pt(size - lvl * 2)
        run.font.color.rgb = color
        p.space_after = Pt(4)
    return tb


def add_rect(slide, left, top, width, height, fill=LIGHT, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.color.rgb = line if line else fill
    s.shadow.inherit = False
    return s


def add_title_bar(slide, title, subtitle=None):
    add_rect(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.9), fill=NAVY)
    add_text(slide, Inches(0.5), Inches(0.18), Inches(12), Inches(0.6),
             title, size=28, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    if subtitle:
        add_text(slide, Inches(0.5), Inches(0.95), Inches(12), Inches(0.4),
                 subtitle, size=14, italic=True, color=MUTED)


# ----- SLIDE 1: TITLE ---------------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_rect(s, Inches(0), Inches(0), prs.slide_width, prs.slide_height, fill=NAVY)
add_text(s, Inches(0.8), Inches(2.4), Inches(11.7), Inches(1.4),
         "Learning-Augmented S3-FIFO",
         size=48, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
add_text(s, Inches(0.8), Inches(3.6), Inches(11.7), Inches(0.8),
         "A small linear gate, an honest result",
         size=24, italic=True, color=RGBColor(0xE0, 0xE0, 0xE0))
add_text(s, Inches(0.8), Inches(5.5), Inches(11.7), Inches(0.5),
         "Minkai Li  ·  CS 2640  ·  Spring 2026",
         size=18, color=RGBColor(0xC8, 0xC8, 0xC8))

# ----- SLIDE 2: GOAL ----------------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_title_bar(s, "Goal: simple learning to augment caching",
              "Not learn-the-policy. Learn one decision inside an existing policy.")

add_text(s, Inches(0.6), Inches(1.3), Inches(12.2), Inches(0.6),
         "Production caches are dominated by hand-crafted heuristics: LRU, FIFO, ARC, S3-FIFO.",
         size=20, color=NAVY)
add_text(s, Inches(0.6), Inches(2.0), Inches(12.2), Inches(0.6),
         "Each heuristic encodes a hardcoded classifier — \"promote on hit,\" \"evict the tail.\"",
         size=20, color=NAVY)

add_rect(s, Inches(0.6), Inches(3.1), Inches(12.1), Inches(2.6), fill=LIGHT, line=NAVY)
add_text(s, Inches(0.9), Inches(3.3), Inches(11.5), Inches(0.5),
         "Our thesis", size=22, bold=True, color=ACCENT)
add_bullets(s, Inches(0.9), Inches(3.85), Inches(11.5), Inches(1.7), [
    "Replace the cheapest, highest-leverage hardcoded decision with a 4-parameter online linear model.",
    "Use only cache-internal signals — no tenant IDs, no content-type embeddings, no offline training.",
    "Keep the surrounding policy structure (S3-FIFO's S→M→G queues) untouched.",
], size=18)

add_rect(s, Inches(0.6), Inches(6.0), Inches(12.1), Inches(1.1),
         fill=RGBColor(0xFF, 0xF6, 0xE0), line=ACCENT)
add_text(s, Inches(0.8), Inches(6.12), Inches(11.7), Inches(0.4),
         "Guest lecture, Valentin Flunkert (AWS)",
         size=14, bold=True, color=ACCENT)
add_text(s, Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.6),
         "No learned cache has shipped at AWS — effectiveness alone isn't the bar; "
         "effectiveness + interpretability is.",
         size=14, italic=True, color=NAVY)

# ----- SLIDE 3: LITERATURE ----------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_title_bar(s, "The template: HALP",
              "Heuristic-Augmented Learned Policies — the deployed model we are imitating.")

# Header card
add_rect(s, Inches(0.4), Inches(1.3), Inches(12.55), Inches(1.3), fill=LIGHT, line=NAVY)
add_text(s, Inches(0.6), Inches(1.45), Inches(12.2), Inches(0.5),
         "Song et al., \"Cache Eviction in Practice: A Heuristic-Augmented Learned Policy,\" NSDI '23",
         size=18, bold=True, color=NAVY)
add_text(s, Inches(0.6), Inches(1.95), Inches(12.2), Inches(0.6),
         "Deployed at YouTube's CDN. 9.1% byte-MR reduction at 1.8% CPU overhead. "
         "First major learned-cache result that survived a real production rollout.",
         size=14, italic=True, color=MUTED)

# How HALP works
add_rect(s, Inches(0.4), Inches(2.85), Inches(12.55), Inches(4.05), fill=LIGHT, line=ACCENT)
add_text(s, Inches(0.6), Inches(2.98), Inches(12.0), Inches(0.5),
         "How HALP works", size=18, bold=True, color=ACCENT)
add_bullets(s, Inches(0.6), Inches(3.5), Inches(12.0), Inches(3.3), [
    "Heuristic (LRU-class) pre-selects k=4 eviction candidates from the tail.",
    "A small continuously-trained MLP scores each candidate.",
    "Pairwise (RLHF-style) preference loss against Belady-derived \"correct\" pairs.",
    "Evict argmin score. Heuristic stays in charge of the search space; learner only re-ranks.",
], size=14)

# ----- SLIDE 5: IDEAS CONSIDERED ----------------------------------------------
s = prs.slides.add_slide(BLANK)
add_title_bar(s, "Learned-caching ideas considered",
              "★ marks the one we pursued.")

ideas = [
    ("★  Learned S → M promotion gate",
     "Replace S3-FIFO's hardcoded promote-on-accessed-bit rule "
     "with a tiny online classifier.",
     ACCENT),
    ("HALP-style rerank",
     "Score k tail candidates at eviction time, evict argmin.",
     NAVY),
    ("Adaptive |S| / |M| sizing",
     "Bandit-tune the S fraction online — Yang et al.'s open problem.",
     NAVY),
    ("Experts meta-policy (vs CACHEUS)",
     "Hedge over {S3-FIFO + gate, vanilla S3-FIFO} or {S3-FIFO, SIEVE}. "
     "Safest \"never lose,\" lowest novelty.",
     NAVY),
    ("Learned admission to S",
     "Refuse to admit large one-hit-wonders into S "
     "(AdaptSize / Flashield style).",
     NAVY),
]
y = Inches(1.45)
row_h = Inches(1.05)
for i, (name, body, c) in enumerate(ideas):
    yy = y + i * row_h
    add_rect(s, Inches(0.4), yy, Inches(12.55), row_h - Inches(0.1),
             fill=LIGHT if i % 2 == 0 else RGBColor(0xEC, 0xEC, 0xEC), line=LIGHT)
    add_text(s, Inches(0.6), yy + Inches(0.2), Inches(4.4), Inches(0.5),
             name, size=18, bold=True, color=c)
    add_text(s, Inches(5.2), yy + Inches(0.25), Inches(7.5), Inches(0.7),
             body, size=14, color=NAVY)

add_text(s, Inches(0.4), Inches(6.85), Inches(12.5), Inches(0.4),
         "Rest of this talk is the promotion gate.",
         size=14, italic=True, color=ACCENT)

# ----- SLIDE 6: WHY THE PROMOTION GATE ----------------------------------------
s = prs.slides.add_slide(BLANK)
add_title_bar(s, "Why the S→M promotion gate?",
              "S3-FIFO's promotion rule is a hardcoded binary classifier — and it's wrong sometimes.")

add_rect(s, Inches(0.4), Inches(1.3), Inches(12.55), Inches(5.6), fill=LIGHT, line=ACCENT)
add_text(s, Inches(0.6), Inches(1.45), Inches(12.0), Inches(0.5),
         "Where the rule fails", size=20, bold=True, color=ACCENT)
add_bullets(s, Inches(0.6), Inches(2.05), Inches(12.0), Inches(4.7), [
    "1-bit summary collapses too much.",
    ("One touch ≡ ten touches; can't separate \"genuinely hot\" from \"barely touched once.\"", 1),
    "Position in S is ignored.",
    ("A touch right after insertion is strong signal; a touch right before eviction is often noise.", 1),
    "Recency of the last hit is invisible.",
    ("Touched 1 ms ago vs 10 min ago = same bit.", 1),
    "Asymmetric error costs.",
    ("False-promote evicts a hot M-resident; false-demote loses a popular newcomer.", 1),
    "Static across workloads.",
    ("Heavy-tailed KV wants \"promote almost nothing\"; locality-rich block I/O wants the opposite.", 1),
], size=15)

# ----- SLIDE 7: THE MODEL -----------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_title_bar(s, "The model: V4",
              "An online logistic regression with 3 features and 4 parameters.")

# Equation
add_rect(s, Inches(0.4), Inches(1.3), Inches(12.55), Inches(1.6), fill=LIGHT, line=NAVY)
add_text(s, Inches(0.6), Inches(1.45), Inches(12.2), Inches(0.5),
         "Decision rule (at the S → M boundary)", size=18, bold=True, color=NAVY)
add_text(s, Inches(0.6), Inches(1.95), Inches(12.2), Inches(0.6),
         "promote(x)  =  σ(b + w₁·log(hits) + w₂·age_S + w₃·recency)  >  0.5",
         size=22, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
add_text(s, Inches(0.6), Inches(2.5), Inches(12.2), Inches(0.4),
         "Online SGD update against the Belady-binary label  y = 1{ next_access_vtime − now < cache_size }",
         size=14, italic=True, color=MUTED, align=PP_ALIGN.CENTER)

# Features
add_text(s, Inches(0.4), Inches(3.2), Inches(12.55), Inches(0.5),
         "Features (each computed from cache-internal state only)",
         size=18, bold=True, color=NAVY)

feats = [
    ("log(hits)",  "log(1 + hits since insertion)"),
    ("age_S",      "(now − insertion_time) / |S|"),
    ("recency",    "(now − last_hit_time) / |S|"),
]
for i, (name, formula) in enumerate(feats):
    yy = Inches(3.85) + i * Inches(0.95)
    add_rect(s, Inches(0.4), yy, Inches(12.55), Inches(0.85),
             fill=RGBColor(0xF8, 0xF8, 0xF8), line=NAVY)
    add_text(s, Inches(0.6), yy + Inches(0.22), Inches(3.5), Inches(0.5),
             name, size=20, bold=True, color=ACCENT)
    add_text(s, Inches(4.5), yy + Inches(0.25), Inches(8.0), Inches(0.5),
             formula, size=15, color=MUTED, italic=True)

# ----- SLIDE 8: ATTEMPTED FEATURES --------------------------------------------
s = prs.slides.add_slide(BLANK)
add_title_bar(s, "Features we tried that didn't help",
              "What got dropped from V3 → V4, and what we tried adding back.")

rows = [
    ("recency / age (ratio)",            "wrong basis"),
    ("log(age)",                         "compresses where decisions live"),
    ("piecewise-linear GAM (16 params)", "undertrained"),
]

add_text(s, Inches(0.4), Inches(1.4), Inches(5.5), Inches(0.5),
         "Feature", size=18, bold=True, color=NAVY)
add_text(s, Inches(7.0), Inches(1.4), Inches(5.5), Inches(0.5),
         "Why it didn't help", size=18, bold=True, color=NAVY)

for i, (f, why) in enumerate(rows):
    yy = Inches(2.0) + i * Inches(1.3)
    add_rect(s, Inches(0.4), yy, Inches(12.55), Inches(1.2),
             fill=LIGHT if i % 2 == 0 else RGBColor(0xEC, 0xEC, 0xEC), line=LIGHT)
    add_text(s, Inches(0.55), yy + Inches(0.4), Inches(6.3), Inches(0.5),
             f, size=18, bold=True, color=ACCENT)
    add_text(s, Inches(7.0), yy + Inches(0.4), Inches(5.5), Inches(0.5),
             why, size=16, italic=True, color=NAVY)

add_text(s, Inches(0.4), Inches(6.45), Inches(12.5), Inches(0.5),
         "Lesson: more expressive ≠ better. The Belady boundary is locally linear in (log_hits, age, recency).",
         size=14, italic=True, color=ACCENT)

# ----- SLIDE 9: SGD IN PRODUCTION ---------------------------------------------
s = prs.slides.add_slide(BLANK)
add_title_bar(s, "How would the SGD actually run in production?",
              "Delayed online updates, inspired by HALP — applied to our promotion gate.")

# Header — the core problem
add_rect(s, Inches(0.4), Inches(1.3), Inches(12.55), Inches(0.95), fill=LIGHT, line=NAVY)
add_text(s, Inches(0.6), Inches(1.42), Inches(12.2), Inches(0.5),
         "Core problem: the Belady label requires knowing the next access — which we don't have at decision time.",
         size=15, bold=True, color=NAVY)
add_text(s, Inches(0.6), Inches(1.83), Inches(12.2), Inches(0.4),
         "Solution: don't train at decision time. Snapshot the decision; train when the label resolves.",
         size=13, italic=True, color=MUTED)

# How it would work for our system
add_rect(s, Inches(0.4), Inches(2.4), Inches(12.55), Inches(4.55), fill=LIGHT, line=ACCENT)
add_text(s, Inches(0.6), Inches(2.5), Inches(12.0), Inches(0.5),
         "Delayed online updates for the S→M gate", size=18, bold=True, color=ACCENT)
add_bullets(s, Inches(0.6), Inches(3.05), Inches(12.0), Inches(3.85), [
    "At each promotion decision, snapshot ⟨ obj_id, features φ(x), gate output ⟩ into a pending-buffer.",
    ("φ(x) = (log_hits, age_S, recency); the snapshot freezes them at decision time.", 1),
    "Wait for the label to resolve. Two natural resolution events:",
    ("Hit  →  the object was accessed again within the horizon W ≈ |cache|. Label y = 1.", 1),
    ("Eviction without a hit  →  the object aged out of M (or G) without a re-access. Label y = 0.", 1),
    "On resolution: apply one SGD step on (φ, y), then drop the tuple from the buffer.",
    "Pending-buffer size is bounded by the cache's own residency window — no unbounded growth.",
    "Failure mode is graceful: if the gate's rolling accuracy degrades, fall back to vanilla S3-FIFO.",
], size=14)

# Bottom — what we did and what we deferred
add_rect(s, Inches(0.4), Inches(7.05), Inches(12.55), Inches(0.4),
         fill=RGBColor(0xFF, 0xF6, 0xE0), line=ACCENT)
add_text(s, Inches(0.55), Inches(7.1), Inches(12.4), Inches(0.35),
         "For this project we use offline Belady labels — an upper bound on what the delayed-online scheme above can learn.",
         size=12, italic=True, color=NAVY)

# ----- SLIDE 10: EARLY RESULTS ------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_title_bar(s, "First results: V4 vs vanilla S3-FIFO (T=1)",
              "On our initial trace set, V4 looked like a clean win. (Spoiler: read on.)")

add_text(s, Inches(0.5), Inches(1.3), Inches(12.3), Inches(0.4),
         "Request miss-ratio, lower is better. 500k requests per cell.",
         size=14, italic=True, color=MUTED)

# Table header
hdrs = ["Trace", "Cache", "S3-FIFO (T=1)", "V4", "Δ"]
xs = [Inches(0.5), Inches(3.4), Inches(5.6), Inches(8.0), Inches(10.5)]
ws = [Inches(2.9), Inches(2.2), Inches(2.4), Inches(2.5), Inches(2.5)]
yh = Inches(1.85)
for x, w, h in zip(xs, ws, hdrs):
    add_rect(s, x, yh, w - Inches(0.05), Inches(0.5), fill=NAVY, line=NAVY)
    add_text(s, x + Inches(0.05), yh + Inches(0.08), w, Inches(0.4),
             h, size=14, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))

# Table rows
rows = [
    ("Twitter cluster52", "1,000", "0.3355", "0.3184", "−1.71 pp ✓"),
    ("Twitter cluster52", "10,000", "0.2110", "0.2078", "−0.32 pp ✓"),
    ("Twitter cluster45 (write-heavy)", "1,000", "0.5630", "0.5459", "−1.71 pp ✓"),
    ("Twitter cluster45", "10,000", "0.4982", "0.4940", "−0.42 pp ✓"),
    ("MSR hm_0 (block I/O)", "1,000", "0.3715", "0.3704", "−0.11 pp ✓"),
    ("CloudPhysics", "5,000", "0.7482", "0.7365", "−1.17 pp ✓"),
]
for i, row in enumerate(rows):
    yy = yh + Inches(0.55) + i * Inches(0.55)
    fill = LIGHT if i % 2 == 0 else RGBColor(0xEC, 0xEC, 0xEC)
    for x, w, val in zip(xs, ws, row):
        add_rect(s, x, yy, w - Inches(0.05), Inches(0.5), fill=fill, line=fill)
        c = GREEN if "✓" in val else NAVY
        add_text(s, x + Inches(0.1), yy + Inches(0.1), w, Inches(0.4),
                 val, size=13, color=c, bold=("✓" in val))

add_text(s, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.5),
         "Headline at this point: V4 wins on 6/6 real cells. Largest single win: −1.71 pp on cluster52 / cluster45.",
         size=15, bold=True, color=GREEN)
add_text(s, Inches(0.5), Inches(6.85), Inches(12.3), Inches(0.4),
         "But: this is V4 vs S3-FIFO at T=1, S=10%. Not the right baseline.",
         size=14, italic=True, color=ACCENT)

# ----- SLIDE 11: THE HONEST RESULT --------------------------------------------
s = prs.slides.add_slide(BLANK)
add_title_bar(s, "The honest result: tuning beats learning",
              "Three same-knob-count ablations against the right baselines.")

# Glossary card
add_rect(s, Inches(0.4), Inches(1.3), Inches(12.55), Inches(1.1), fill=LIGHT, line=NAVY)
add_text(s, Inches(0.6), Inches(1.4), Inches(12.2), Inches(0.4),
         "Baselines", size=14, bold=True, color=NAVY)
add_text(s, Inches(0.6), Inches(1.75), Inches(12.2), Inches(0.6),
         "OptS = best of S3-FIFO with S ∈ {0.01, 0.05, 0.10, 0.25} at T=1.   "
         "OptST = best of S3-FIFO with S ∈ {0.01, 0.05, 0.10, 0.25} × T ∈ {1, 2, 3}.   "
         "All numbers per-cell, post-hoc.",
         size=12, italic=True, color=MUTED)

# Win/loss tally
add_rect(s, Inches(0.4), Inches(2.55), Inches(12.55), Inches(4.05), fill=LIGHT, line=NAVY)
add_text(s, Inches(0.6), Inches(2.7), Inches(12.2), Inches(0.5),
         "Cross-trace tally — 14 traces, 28 effective cells",
         size=16, bold=True, color=NAVY)

# Header row
hdr_y = Inches(3.4)
add_text(s, Inches(0.6),  hdr_y, Inches(7.5), Inches(0.4), "Comparison", size=13, bold=True, color=MUTED)
add_text(s, Inches(8.4),  hdr_y, Inches(1.2), Inches(0.4), "Wins",   size=13, bold=True, color=MUTED, align=PP_ALIGN.CENTER)
add_text(s, Inches(9.7),  hdr_y, Inches(1.2), Inches(0.4), "Losses", size=13, bold=True, color=MUTED, align=PP_ALIGN.CENTER)
add_text(s, Inches(11.0), hdr_y, Inches(1.2), Inches(0.4), "Ties",   size=13, bold=True, color=MUTED, align=PP_ALIGN.CENTER)

tallies = [
    ("V4  vs  OptST           (gate vs full-grid tuning)",            "3",  "16", "9"),
    ("V4 + OptS  vs  OptST    (gate + S sweep vs full-grid tuning)",  "5",  "10", "13"),
    ("V4 + OptS  vs  OptS     (same-knob-count ablation)",            "8",  "7",  "13"),
]
for i, (label, w, l, t) in enumerate(tallies):
    yy = Inches(4.05) + i * Inches(0.75)
    add_text(s, Inches(0.6),  yy, Inches(7.6), Inches(0.6),
             label, size=15, color=NAVY)
    add_text(s, Inches(8.4),  yy, Inches(1.2), Inches(0.6), w, size=20, bold=True, color=GREEN,  align=PP_ALIGN.CENTER)
    add_text(s, Inches(9.7),  yy, Inches(1.2), Inches(0.6), l, size=20, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    add_text(s, Inches(11.0), yy, Inches(1.2), Inches(0.6), t, size=20, bold=True, color=MUTED,  align=PP_ALIGN.CENTER)

# Trace inventory — small italic footer
add_text(s, Inches(0.4), Inches(6.8), Inches(12.5), Inches(0.35),
         "Traces (14): Twitter cluster26, 45, 50, 52  ·  MSR hm_0, proj_0, prxy_0  ·  "
         "Alibaba 110, 185  ·  CloudPhysics, w105  ·  Wiki, meta_reag  ·  Meta block1.  "
         "(cluster10 excluded — zero signal.)",
         size=10, italic=True, color=MUTED)
add_text(s, Inches(0.4), Inches(7.15), Inches(12.5), Inches(0.35),
         "14 traces × 2 cache sizes = 28 effective cells. 500k requests per cell where applicable.",
         size=10, italic=True, color=MUTED)

# ----- SLIDE 12: WHEN LEARNING HELPS / FAILS (CONSOLIDATED) -------------------
s = prs.slides.add_slide(BLANK)
add_title_bar(s, "Where learned caching matters — and where it doesn't")

# ============= LEFT — When it helps ============================================
add_text(s, Inches(0.4), Inches(1.15), Inches(6.2), Inches(0.5),
         "✔  When learning helps", size=20, bold=True, color=GREEN)

add_bullets(s, Inches(0.4), Inches(1.8), Inches(6.2), Inches(5.5), [
    "The right policy is conditional — \"promote when X is true\" — and a static knob can't express the condition.",
    ("CloudPhysics cache=5,000  (Δ −1.72 pp vs OptS)", 1),
    ("w_recency = −1.08 dominates;  OptS 0.7505 → V4+OptS 0.7333.", 1),
    "Heavy-tailed workload where \"promote almost nothing\" is right, and the gate auto-discovers it.",
    ("Twitter cluster52 cache=1,000  (Δ −1.01 pp vs OptS)", 1),
    ("Gate converges to z_max ≈ −1.8;  OptS 0.3234 → V4+OptS 0.3133.", 1),
    "Common thread: the trace's reuse structure has variance the W-windowed Belady label can capture.",
], size=14, color=NAVY)

# ============= RIGHT — When it fails ===========================================
add_text(s, Inches(6.75), Inches(1.15), Inches(6.2), Inches(0.5),
         "✘  When learning fails", size=20, bold=True, color=ACCENT)

add_bullets(s, Inches(6.75), Inches(1.8), Inches(6.2), Inches(3.0), [
    "The W-windowed Belady label cuts through a sharp peak in the trace's reuse-distance distribution.",
    ("msr_prxy_0 cache=1,000  (V4+OptS +8.23 pp WORSE)", 1),
    ("77% of reuses at distance 1024–2047; W=1,000 calls them all y=0.", 1),
    ("Gate demotes them; they all come back at ~1,500. Verified by histogram.", 1),
], size=14, color=NAVY)

# Speculative footer
add_rect(s, Inches(6.75), Inches(4.95), Inches(6.2), Inches(2.4),
         fill=RGBColor(0xF8, 0xF8, 0xF8), line=MUTED)
add_text(s, Inches(6.9), Inches(5.05), Inches(5.9), Inches(0.4),
         "Other speculative mechanisms (no clean trace example):",
         size=12, bold=True, italic=True, color=MUTED)
add_bullets(s, Inches(6.9), Inches(5.5), Inches(5.9), Inches(1.85), [
    "Permissive labels at large cache — almost everything has y=1, gate has no negative examples to learn from.",
    "Learning too aggressive or too conservative — weight magnitudes pushing decisions to extremes (or near the bias).",
    "Coefficient noise / non-convergence — SGD doesn't settle within trace length; weights drift.",
    "Off-policy data distribution — gate trains on data its own past decisions admitted.",
], size=11, color=MUTED)

# ----- SLIDE 13: NEXT STEPS ---------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_title_bar(s, "Next steps")

next_steps = [
    ("Diagnose the failure modes",
     "measure y=1 fraction at decision time   ·   track SGD weight trajectories   ·   "
     "alternative label horizons (W = 0.5·cache, 2·cache)   ·   pairwise loss à la HALP"),
    ("Add richer, cache-wide features",
     "lifetime hits   ·   ghost-hit count   ·   rolling MR over last K requests   ·   "
     "current S occupancy   ·   age-distribution percentile of x   ·   log(size) on real-size traces"),
    ("Apply the gate at more decision points",
     "M-eviction gate (same model, second binary classifier inside S3-FIFO; ~50 LoC)   ·   "
     "admission to S   ·   mirror-descent |S| sizer to replace the EXP3 attempt that didn't converge"),
]
y0 = Inches(1.5)
row_h = Inches(1.65)
for i, (title, sub) in enumerate(next_steps):
    yy = y0 + i * row_h
    add_rect(s, Inches(0.4), yy, Inches(12.55), row_h - Inches(0.2), fill=LIGHT, line=NAVY)
    add_text(s, Inches(0.7), yy + Inches(0.2), Inches(12.0), Inches(0.5),
             "•  " + title, size=20, bold=True, color=NAVY)
    add_text(s, Inches(1.05), yy + Inches(0.8), Inches(11.7), Inches(0.65),
             "–  " + sub, size=14, color=MUTED, italic=True)

add_text(s, Inches(0.4), Inches(6.85), Inches(12.5), Inches(0.4),
         "Highest expected lift: M-eviction gate.   Most diagnostic value: actual decision-time label distributions.",
         size=13, italic=True, color=ACCENT)

# ----- SLIDE 16: THANKS -------------------------------------------------------
s = prs.slides.add_slide(BLANK)
add_rect(s, Inches(0), Inches(0), prs.slide_width, prs.slide_height, fill=NAVY)
add_text(s, Inches(0.8), Inches(2.6), Inches(11.7), Inches(1.3),
         "Thanks — questions?",
         size=44, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
add_text(s, Inches(0.8), Inches(4.0), Inches(11.7), Inches(1.4),
         ["Code:  plugins/learned_promotion.py  (V4 family)",
          "Sweep:  plugins/sweep14.py  ·  plugins/sweep14_report.py",
          "Numbers come from libCacheSim's canonical baselines + a Python harness for the gate."],
         size=16, color=RGBColor(0xE0, 0xE0, 0xE0))

prs.save(OUT)
print(f"Saved {OUT} ({len(prs.slides)} slides)")
